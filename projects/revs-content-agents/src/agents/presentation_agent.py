from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from src.parsing.schemas import DesignBrief
from src.utils.config import load_config


def _hex_to_rgb(hex_color: str):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(hex_color.lstrip("#"))


class PresentationAgent:
    format_name = "presentations"

    def generate(self, brief: DesignBrief, output_dir: Path) -> dict:
        out_dir = output_dir / self.format_name
        out_dir.mkdir(parents=True, exist_ok=True)
        branding = load_config().get("branding", {})
        primary = _hex_to_rgb(branding.get("primary_color", "#0052CC"))

        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        body_layout = prs.slide_layouts[1]

        title_slide = prs.slides.add_slide(title_layout)
        title_slide.shapes.title.text = brief.headline
        title_slide.placeholders[1].text = brief.intro[:280]
        title_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary

        for section in brief.sections:
            slide = prs.slides.add_slide(body_layout)
            slide.shapes.title.text = section.title
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary
            body = slide.placeholders[1].text_frame
            body.word_wrap = True
            paragraphs = [p.strip() for p in section.body.split("\n\n") if p.strip()][:4]
            body.text = paragraphs[0] if paragraphs else ""
            for para in paragraphs[1:]:
                p = body.add_paragraph()
                p.text = para
                p.font.size = Pt(14)
            slide.notes_slide.notes_text_frame.text = section.body

        if brief.closing_title:
            closing_slide = prs.slides.add_slide(body_layout)
            closing_slide.shapes.title.text = brief.closing_title
            closing_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary
            closing_slide.placeholders[1].text = brief.closing_body or ""

        pptx_path = out_dir / f"{brief.concept_id}.pptx"
        prs.save(str(pptx_path))

        return {"format": self.format_name, "files": [str(pptx_path)]}
