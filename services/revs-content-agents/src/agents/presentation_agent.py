import re
import textwrap
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

from src.parsing.schemas import DesignBrief
from src.utils.config import load_config

_HEADING_RE = re.compile(r"^#{1,6}\s*")
_BOLD_SPLIT_RE = re.compile(r"(\*\*.*?\*\*)")


def _hex_to_rgb(hex_color: str):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(hex_color.lstrip("#"))


def _add_markdown_paragraph(text_frame, text: str, is_first: bool, font_size: int | None = None) -> None:
    """Render `**bold**` as real bold runs and strip leading '#' heading markers,
    instead of dumping raw Markdown syntax as literal slide text."""
    text = _HEADING_RE.sub("", text.strip(), count=1)
    paragraph = text_frame.paragraphs[0] if is_first else text_frame.add_paragraph()
    for part in _BOLD_SPLIT_RE.split(text):
        if not part:
            continue
        run = paragraph.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]
            run.font.bold = True
        else:
            run.text = part
        if font_size:
            run.font.size = Pt(font_size)


class PresentationAgent:
    format_name = "presentations"

    def generate(self, brief: DesignBrief, output_dir: Path) -> dict:
        out_dir = output_dir / self.format_name
        out_dir.mkdir(parents=True, exist_ok=True)
        branding = load_config().get("branding", {})
        primary = _hex_to_rgb(branding.get("primary_color", "#0052CC"))

        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        body_layout = prs.slide_layouts[1]

        title_slide = prs.slides.add_slide(title_layout)
        title_slide.shapes.title.text = brief.headline
        title_slide.placeholders[1].text = textwrap.shorten(brief.intro, width=280, placeholder="...")
        title_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary

        for section in brief.sections:
            slide = prs.slides.add_slide(body_layout)
            slide.shapes.title.text = section.title
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary
            body = slide.placeholders[1].text_frame
            body.word_wrap = True
            paragraphs = [p.strip() for p in section.body.split("\n\n") if p.strip()]
            for i, para in enumerate(paragraphs):
                _add_markdown_paragraph(body, para, is_first=(i == 0), font_size=14 if i > 0 else None)
            slide.notes_slide.notes_text_frame.text = section.body

        if brief.closing_title:
            closing_slide = prs.slides.add_slide(body_layout)
            closing_slide.shapes.title.text = brief.closing_title
            closing_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary
            closing_body = closing_slide.placeholders[1].text_frame
            closing_paragraphs = [p.strip() for p in (brief.closing_body or "").split("\n\n") if p.strip()]
            for i, para in enumerate(closing_paragraphs):
                _add_markdown_paragraph(closing_body, para, is_first=(i == 0))

        pptx_path = out_dir / f"{brief.concept_id}.pptx"
        prs.save(str(pptx_path))

        return {"format": self.format_name, "files": [str(pptx_path)]}
