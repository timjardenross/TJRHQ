from pptx import Presentation

from src.agents.presentation_agent import PresentationAgent


def test_generate_presentation_has_title_plus_section_plus_closing_slides(brief, tmp_path):
    result = PresentationAgent().generate(brief, tmp_path)
    prs = Presentation(result["files"][0])

    # 1 title slide + 5 section slides + 1 closing slide
    assert len(prs.slides._sldIdLst) == 7
    assert prs.slides[0].shapes.title.text == brief.headline
    assert prs.slides[-1].shapes.title.text == brief.closing_title
